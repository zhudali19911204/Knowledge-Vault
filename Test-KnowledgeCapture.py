"""Capture regressions using synthetic documents, never a user's source files."""

import base64
import importlib.util
import io
from pathlib import Path
import shutil
import sys
import tempfile
import unittest
from unittest.mock import patch

HAS_PPTX = importlib.util.find_spec("pptx") is not None
if HAS_PPTX:
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    from pptx.util import Inches


SCRIPT = Path(__file__).resolve().parent / "vault-template/.dsh/skills/knowledge-capture/scripts/document_to_markdown.py"
spec = importlib.util.spec_from_file_location("capture_converter", SCRIPT)
converter = importlib.util.module_from_spec(spec)
sys.modules[spec.name] = converter
spec.loader.exec_module(converter)
PNG = base64.b64decode("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+j0f8AAAAASUVORK5CYII=")


class CaptureTestCase(unittest.TestCase):
    def setUp(self):
        self.temp = tempfile.TemporaryDirectory(prefix="knowledge-capture-test-")
        self.root = Path(self.temp.name).resolve()
        self.assertEqual(self.root.parent, Path(tempfile.gettempdir()).resolve())
        self.assertTrue(self.root.name.startswith("knowledge-capture-test-"))
        self.addCleanup(self.temp.cleanup)
        self.vault = self.root / "vault"
        self.vault.mkdir()
        (self.vault / "AGENTS.md").write_text("# Test Vault\n", encoding="utf-8")
        for name in ("01_Inbox", "07_Attachments"):
            (self.vault / name).mkdir()

    def run_conversion(self, source, mode="attachments"):
        before_hash = converter.sha256_file(source)
        before_stat = source.stat()
        result = converter.run_conversion(source, self.vault, mode, None)
        self.assertEqual(result["status"], "ok")
        self.assertEqual(before_hash, converter.sha256_file(source))
        self.assertEqual(before_stat.st_mtime_ns, source.stat().st_mtime_ns)
        return result, Path(result["markdown"]).read_text(encoding="utf-8")


@unittest.skipUnless(HAS_PPTX, "python-pptx unavailable; run with the installed knowledge-capture runtime Python")
class PowerPointCaptureTests(CaptureTestCase):
    def presentation(self, missing="linked", valid=True, group=False):
        document = Presentation()
        slide = document.slides.add_slide(document.slide_layouts[6])
        shapes = slide.shapes.add_group_shape().shapes if group else slide.shapes
        picture = shapes.add_picture(io.BytesIO(PNG), Inches(1), Inches(1), width=Inches(1))
        picture.name = "missing-picture"
        blip = picture._element.blipFill.blip
        if missing == "linked":
            del blip.attrib[qn("r:embed")]
            blip.set(qn("r:link"), slide.part.relate_to("https://example.invalid/private-image.png", RT.IMAGE, is_external=True))
        elif missing == "empty":
            del blip.attrib[qn("r:embed")]
        elif missing == "broken":
            blip.set(qn("r:embed"), "rId-does-not-exist")
        if valid:
            shapes.add_picture(io.BytesIO(PNG), Inches(1), Inches(3), width=Inches(1))
        slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1)).text = "Text after the unavailable picture."
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(6), Inches(5), Inches(1)).table
        table.cell(0, 0).text, table.cell(0, 1).text = "Key", "Value"
        slide.notes_slide.notes_text_frame.text = "Keep speaker notes."
        second = document.slides.add_slide(document.slide_layouts[6])
        second.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Slide two stays present."
        source = self.root / "input.pptx"
        document.save(source)
        return source

    def test_linked_picture_does_not_abort_text_table_notes_or_other_images(self):
        result, text = self.run_conversion(self.presentation())
        self.assertEqual(result["images_processed"], 1)
        self.assertEqual(result["attachments_written"], 1)
        self.assertEqual(result["details"]["images_unavailable"], 1)
        self.assertIn("外部链接", text)
        self.assertIn("幻灯片 1·对象 1", text)
        self.assertIn("Text after the unavailable picture.", text)
        self.assertIn("| Key | Value |", text)
        self.assertIn("Keep speaker notes.", text)
        self.assertIn("Slide two stays present.", text)
        self.assertIn("## 转换说明", text)
        self.assertEqual(text.count("![[07_Attachments/"), 1)
        self.assertEqual(len(list(Path(result["attachment_directory"]).iterdir())), 1)
        self.assertNotIn("https://example.invalid", text)

    def test_empty_or_broken_picture_reference_is_reported_at_its_position(self):
        for missing in ("empty", "broken"):
            with self.subTest(missing=missing):
                result, text = self.run_conversion(self.presentation(missing))
                self.assertEqual(result["details"]["images_unavailable"], 1)
                self.assertIn("图片未能提取", text)
                self.assertLess(text.index("图片未能提取"), text.index("Text after"))
                self.assertEqual(result["attachments_written"], 1)

    def test_grouped_picture_failure_is_isolated(self):
        result, text = self.run_conversion(self.presentation(group=True))
        self.assertEqual(result["details"]["images_unavailable"], 1)
        self.assertEqual(result["attachments_written"], 1)
        self.assertIn("Slide two stays present.", text)

    def test_filled_picture_placeholder_is_extracted(self):
        document = Presentation()
        slide = document.slides.add_slide(document.slide_layouts[8])
        slide.placeholders[1].insert_picture(io.BytesIO(PNG))
        source = self.root / "placeholder.pptx"
        document.save(source)
        result, text = self.run_conversion(source)
        self.assertEqual(result["attachments_written"], 1)
        self.assertEqual(result["details"]["images_unavailable"], 0)
        self.assertIn("![[07_Attachments/", text)

    def test_multimodal_with_only_missing_images_finishes_with_warnings(self):
        source = self.presentation(valid=False)
        result = converter.prepare_multimodal(source, self.vault, None)
        self.assertEqual(result["status"], "ok")
        self.assertEqual(result["images_processed"], 0)
        self.assertEqual(result["details"]["images_unavailable"], 1)
        self.assertIsNone(result["attachment_directory"])
        text = Path(result["markdown"]).read_text(encoding="utf-8")
        self.assertIn("图片未能提取", text)
        self.assertIn("Slide two stays present.", text)

    def test_multimodal_keeps_missing_image_warning_and_stages_valid_picture(self):
        source = self.presentation()
        result = converter.prepare_multimodal(source, self.vault, None)
        self.assertEqual(result["status"], "needs_multimodal")
        manifest_path = Path(result["manifest"])
        job_root = converter.validate_multimodal_job_path(manifest_path)
        self.addCleanup(shutil.rmtree, job_root)
        self.assertEqual(result["images_to_read"], 1)
        self.assertEqual(result["details"]["images_unavailable"], 1)
        self.assertTrue(result["warnings"])
        self.assertIn("图片未能提取", (job_root / "draft.md").read_text(encoding="utf-8"))
        self.assertFalse(list((self.vault / "01_Inbox").iterdir()))

    def test_attachment_storage_failure_is_not_hidden_as_a_missing_picture(self):
        source = self.presentation(missing=None)
        with patch.object(converter.ImageStager, "add", side_effect=PermissionError("storage denied")):
            with self.assertRaises(PermissionError):
                converter.run_conversion(source, self.vault, "attachments", None)
        self.assertFalse(list((self.vault / "01_Inbox").iterdir()))


class TextCaptureTests(CaptureTestCase):
    def test_plain_csv_commit_still_writes_complete_note(self):
        source = self.root / "input.csv"
        source.write_text("name,value\nfirst,42\n", encoding="utf-8")
        result, text = self.run_conversion(source, "text")
        self.assertEqual(result["attachments_written"], 0)
        self.assertIn("first | 42", text)


if __name__ == "__main__":
    unittest.main(verbosity=2)
