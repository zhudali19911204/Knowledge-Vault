#!/usr/bin/env python3
"""Convert Office/PDF sources into a traceable Knowledge Vault Inbox note."""

from __future__ import annotations

import argparse
import csv
import hashlib
import importlib.util
import io
import json
import os
import posixpath
import re
import shutil
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")


SUPPORTED = {".xlsx", ".xlsm", ".csv", ".tsv", ".pdf", ".docx", ".pptx"}
LEGACY = {".xls", ".doc", ".ppt"}
MIXED_LAYOUT = {".pdf", ".docx", ".pptx"}
INVALID_FILENAME = re.compile(r'[\\/:*?"<>|]+')
MULTIMODAL_SCHEMA_VERSION = "knowledge-capture-multimodal/v1"
MULTIMODAL_PLACEHOLDER = "knowledge-capture-multimodal"
MODEL_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
NS = {"w": W_NS, "r": R_NS, "a": A_NS}


class ConversionError(RuntimeError):
    pass


class DependencyError(ConversionError):
    pass


class SourceError(ConversionError):
    def __init__(self, source: Path, error_code: str, message: str) -> None:
        super().__init__(message)
        self.source = source
        self.error_code = error_code


def source_access_error(source: Path, error: OSError) -> SourceError:
    if isinstance(error, FileNotFoundError):
        return SourceError(source, "source_not_found", f"源文件不存在，请核对完整路径：{source}")
    return SourceError(
        source,
        "source_unreadable",
        f"无法读取源文件，尚未开始解析内容：{source}。"
        "请先关闭占用该文件的 Office/WPS 或预览窗口；若仍无法读取，再检查读取权限及 OneDrive 本地可用状态。"
        f"系统信息：{error}",
    )


@dataclass
class ImageRecord:
    identifier: str
    source: Path
    filename: str
    vault_path: str
    location: str
    model_readable: bool = True
    unreadable_reason: str | None = None


@dataclass
class Conversion:
    markdown: str
    warnings: list[str] = field(default_factory=list)
    details: dict = field(default_factory=dict)


def module_available(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def safe_name(value: str, fallback: str = "文档", limit: int = 90) -> str:
    cleaned = INVALID_FILENAME.sub("-", str(value)).strip().strip(".")
    cleaned = re.sub(r"\s+", " ", cleaned)
    return (cleaned or fallback)[:limit].rstrip()


def yaml_string(value: object) -> str:
    return json.dumps(str(value), ensure_ascii=False)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def markdown_cell(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, datetime):
        text = value.isoformat(sep=" ")
    else:
        text = str(value)
    return text.replace("\\", "\\\\").replace("|", "\\|").replace("\r\n", "<br>").replace("\n", "<br>")


def column_name(index: int) -> str:
    result = ""
    while index:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def markdown_table(rows: list[list[object]], include_coordinates: bool = False) -> str:
    width = max((len(row) for row in rows), default=0)
    if width == 0:
        return "_空表格_"
    header = (["行"] if include_coordinates else []) + [column_name(index) for index in range(1, width + 1)]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row_index, row in enumerate(rows, 1):
        values = ([row_index] if include_coordinates else []) + list(row) + [""] * (width - len(row))
        lines.append("| " + " | ".join(markdown_cell(value) for value in values) + " |")
    return "\n".join(lines)


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    for index in range(2, 10000):
        candidate = path.with_name(f"{path.stem} ({index}){path.suffix}")
        if not candidate.exists():
            return candidate
    raise ConversionError(f"无法为输出分配不冲突的文件名：{path}")


def find_attachment_folder(vault_root: Path, title: str) -> tuple[str, Path]:
    root = vault_root / "07_Attachments"
    suffix = safe_name(title)
    existing = sorted(root.glob(f"[0-9][0-9][0-9][0-9]_{suffix}"))
    if existing:
        folder = existing[0]
        return folder.name, folder
    highest = 0
    for entry in root.iterdir():
        if entry.is_dir() and (match := re.match(r"^07(\d{2})_", entry.name)):
            highest = max(highest, int(match.group(1)))
    if highest >= 99:
        raise ConversionError("07_Attachments 下的两位目录编号已用尽。")
    name = f"07{highest + 1:02d}_{suffix}"
    return name, root / name


def inspect_package_images(source: Path, prefix: str) -> int:
    try:
        with zipfile.ZipFile(source) as package:
            return sum(1 for name in package.namelist() if name.startswith(prefix) and not name.endswith("/"))
    except OSError as error:
        raise source_access_error(source, error) from error
    except zipfile.BadZipFile as error:
        raise SourceError(
            source,
            "invalid_document_package",
            f"源文件可以读取，但不是有效的 {source.suffix} 文档包：{source}。"
            "请用原应用确认文件格式及是否加密或损坏；不能仅修改扩展名。",
        ) from error


def inspect_source(source: Path) -> dict:
    extension = source.suffix.lower()
    try:
        source_stat = source.stat()
        source_hash = sha256_file(source)
    except OSError as error:
        raise source_access_error(source, error) from error
    result = {
        "source": str(source),
        "format": extension.lstrip("."),
        "bytes": source_stat.st_size,
        "modified_at": datetime.fromtimestamp(source_stat.st_mtime, timezone.utc).astimezone().isoformat(),
        "sha256": source_hash,
        "image_count": 0,
        "requires_image_mode": False,
        "missing_dependencies": [],
    }
    if extension == ".docx":
        result["image_count"] = inspect_package_images(source, "word/media/")
    elif extension == ".pptx":
        result["image_count"] = inspect_package_images(source, "ppt/media/")
    elif extension in {".xlsx", ".xlsm"}:
        result["image_count"] = inspect_package_images(source, "xl/media/")
        if not module_available("openpyxl"):
            result["missing_dependencies"].append("openpyxl")
    elif extension == ".pdf":
        if not module_available("fitz"):
            result["image_count"] = None
            result["missing_dependencies"].append("PyMuPDF")
        else:
            import fitz

            with fitz.open(source) as document:
                result["image_count"] = sum(len(page.get_images(full=True)) for page in document)
                result["vector_pages"] = [
                    number for number, page in enumerate(document, 1)
                    if pdf_vector_bounds(page) is not None
                ]
                if result["image_count"] == 0:
                    result["image_count"] = sum(1 for page in document if not page.get_text().strip())
    if extension == ".pptx" and not module_available("pptx"):
        result["missing_dependencies"].append("python-pptx")
    result["requires_image_mode"] = extension in MIXED_LAYOUT and (
        result["image_count"] is None or int(result["image_count"] or 0) > 0 or bool(result.get("vector_pages"))
    )
    return result


class ImageStager:
    def __init__(
        self,
        temp_root: Path,
        destination: Path,
        destination_name: str,
        mode: str,
    ) -> None:
        self.temp_root = temp_root
        self.destination = destination
        self.destination_name = destination_name
        self.mode = mode
        self.records: list[ImageRecord] = []
        self.used_names: set[str] = set()
        self.warnings: list[str] = []

    def _allocate_name(self, label: str, extension: str) -> str:
        base = safe_name(label, f"图片-{len(self.records) + 1:03d}", 70)
        extension = extension.lower()
        if not extension.startswith("."):
            extension = "." + extension
        if extension == ".jpeg":
            extension = ".jpg"
        candidate = f"{len(self.records) + 1:03d}_{base}{extension}"
        index = 2
        while candidate.lower() in self.used_names or (self.destination / candidate).exists():
            candidate = f"{len(self.records) + 1:03d}_{base}-{index}{extension}"
            index += 1
        self.used_names.add(candidate.lower())
        return candidate

    def _model_image(self, data: bytes, extension: str) -> tuple[bytes, str, bool, str | None]:
        normalized = extension.lower()
        if not normalized.startswith("."):
            normalized = "." + normalized
        if normalized == ".jpeg":
            normalized = ".jpg"
        if normalized in MODEL_IMAGE_EXTENSIONS:
            return data, normalized, True, None
        if not module_available("PIL"):
            return data, normalized, False, "缺少 Pillow，无法把该图片转换为模型支持的格式"
        try:
            from PIL import Image

            with Image.open(io.BytesIO(data)) as image:
                image.seek(0)
                converted = image.convert("RGBA" if "A" in image.getbands() else "RGB")
                output = io.BytesIO()
                converted.save(output, format="PNG")
            return output.getvalue(), ".png", True, None
        except Exception as error:
            return data, normalized, False, f"无法转换为 PNG：{error}"

    def add(self, data: bytes, extension: str, label: str, location: str) -> str:
        model_readable = True
        unreadable_reason = None
        if self.mode == "multimodal":
            data, extension, model_readable, unreadable_reason = self._model_image(data, extension)
        filename = self._allocate_name(label, extension)
        staged = self.temp_root / filename
        staged.write_bytes(data)
        vault_path = f"07_Attachments/{self.destination_name}/{filename}"
        identifier = f"I{len(self.records) + 1:03d}"
        record = ImageRecord(identifier, staged, filename, vault_path, location, model_readable, unreadable_reason)
        self.records.append(record)
        if self.mode == "multimodal":
            if not model_readable:
                self.warnings.append(f"{location} 的图片无法交给多模态模型：{unreadable_reason}。")
            return f"<!-- {MULTIMODAL_PLACEHOLDER}:{identifier} -->"
        return f"![[{vault_path}]]\n\n_图片位置：{location}_"


def convert_delimited(source: Path, delimiter: str) -> Conversion:
    warnings: list[str] = []
    encoding = "utf-8-sig"
    try:
        text = source.read_text(encoding=encoding)
    except UnicodeDecodeError:
        encoding = "gb18030"
        text = source.read_text(encoding=encoding)
        warnings.append("源文件不是 UTF-8，已按 GB18030 解码。")
    rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
    body = "## 数据表\n\n" + markdown_table(rows, include_coordinates=True)
    return Conversion(body, warnings, {"encoding": encoding, "rows": len(rows), "columns": max(map(len, rows), default=0)})


def convert_excel(source: Path, stager: ImageStager) -> Conversion:
    if not module_available("openpyxl"):
        raise DependencyError("Excel 转换需要 openpyxl；请按 scripts/requirements.txt 安装。")
    import openpyxl

    workbook = openpyxl.load_workbook(source, data_only=False, read_only=False, keep_vba=source.suffix.lower() == ".xlsm")
    cached = openpyxl.load_workbook(source, data_only=True, read_only=False, keep_vba=source.suffix.lower() == ".xlsm")
    sections: list[str] = []
    warnings: list[str] = []
    total_formula_cells = 0
    for sheet in workbook.worksheets:
        cached_sheet = cached[sheet.title]
        rows: list[list[object]] = []
        comments: list[str] = []
        formulas_without_cached_values: list[str] = []
        for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, min_col=1, max_col=sheet.max_column):
            values: list[object] = []
            for cell in row:
                value = cached_sheet[cell.coordinate].value
                if cell.data_type == "f" or (isinstance(cell.value, str) and cell.value.startswith("=")):
                    total_formula_cells += 1
                    if value is None:
                        formulas_without_cached_values.append(cell.coordinate)
                values.append(value)
                if cell.comment is not None:
                    comments.append(f"{cell.coordinate}: {cell.comment.text}")
            rows.append(values)
        hidden_rows = [str(index) for index, dimension in sheet.row_dimensions.items() if dimension.hidden]
        hidden_columns = [str(name) for name, dimension in sheet.column_dimensions.items() if dimension.hidden]
        metadata = [
            f"- 状态：{sheet.sheet_state}",
            f"- 使用范围：A1:{column_name(sheet.max_column)}{sheet.max_row}",
            f"- 冻结窗格：{sheet.freeze_panes or '无'}",
            f"- 筛选范围：{sheet.auto_filter.ref or '无'}",
            f"- 合并区域：{', '.join(str(item) for item in sheet.merged_cells.ranges) or '无'}",
            f"- 隐藏行：{', '.join(hidden_rows) or '无'}",
            f"- 隐藏列：{', '.join(hidden_columns) or '无'}",
        ]
        content = [f"## 工作表：{sheet.title}", "", *metadata, "", markdown_table(rows, include_coordinates=True)]
        if comments:
            content.extend(["", "### 批注", "", *[f"- {value}" for value in comments]])
        images = list(getattr(sheet, "_images", []))
        if images:
            content.extend(["", "### 工作表图片", ""])
            for index, image in enumerate(images, 1):
                anchor = getattr(image, "anchor", None)
                if hasattr(anchor, "_from"):
                    location = f"{sheet.title}!{column_name(anchor._from.col + 1)}{anchor._from.row + 1}"
                else:
                    location = f"{sheet.title}（锚点未知）"
                extension = "." + (getattr(image, "format", None) or "png")
                content.append(stager.add(image._data(), extension, f"{sheet.title}-图片-{index}", location))
        charts = list(getattr(sheet, "_charts", []))
        if charts:
            warnings.append(f"工作表“{sheet.title}”含 {len(charts)} 个图表对象；已保留数据，图表样式无法完整转成 Markdown。")
        if getattr(sheet.data_validations, "count", 0):
            warnings.append(f"工作表“{sheet.title}”含数据验证规则，已记录单元格值但未复刻交互规则。")
        if formulas_without_cached_values:
            warnings.append(
                f"工作表“{sheet.title}”有 {len(formulas_without_cached_values)} 个公式单元格没有缓存实际值，"
                f"Markdown 已留空且未重算：{', '.join(formulas_without_cached_values)}"
            )
        sections.append("\n".join(content))
    if source.suffix.lower() == ".xlsm":
        warnings.append("源文件包含宏容器；转换过程中未执行宏，Markdown 不保留 VBA 行为。")
    return Conversion(
        "\n\n".join(sections),
        warnings,
        {
            "sheets": workbook.sheetnames,
            "formula_cells": total_formula_cells,
        },
    )


def relationship_map(package: zipfile.ZipFile, part_name: str) -> dict[str, str]:
    part = Path(part_name)
    rels_name = posixpath.join(str(part.parent).replace("\\", "/"), "_rels", part.name + ".rels")
    if rels_name not in package.namelist():
        return {}
    root = ET.fromstring(package.read(rels_name))
    result: dict[str, str] = {}
    base = str(part.parent).replace("\\", "/")
    for relation in root.findall(f"{{{PKG_REL_NS}}}Relationship"):
        relation_id = relation.attrib.get("Id")
        target = relation.attrib.get("Target")
        if relation_id and target and not target.startswith(("http:", "https:")):
            result[relation_id] = posixpath.normpath(posixpath.join(base, target))
    return result


def word_paragraph(element: ET.Element, relationships: dict[str, str], package: zipfile.ZipFile, stager: ImageStager, location: str) -> str:
    fragments: list[str] = []
    images: list[str] = []
    image_index = 0
    for node in element.iter():
        if node.tag == f"{{{W_NS}}}t":
            fragments.append(node.text or "")
        elif node.tag == f"{{{W_NS}}}tab":
            fragments.append("\t")
        elif node.tag == f"{{{W_NS}}}br":
            fragments.append("\n")
        elif node.tag == f"{{{A_NS}}}blip":
            relation_id = node.attrib.get(f"{{{R_NS}}}embed")
            target = relationships.get(relation_id or "")
            if target and target in package.namelist():
                image_index += 1
                extension = Path(target).suffix or ".png"
                images.append(stager.add(package.read(target), extension, f"Word-图片-{len(stager.records) + 1}", location))
    text = "".join(fragments).strip()
    style = element.find("./w:pPr/w:pStyle", NS)
    style_name = style.attrib.get(f"{{{W_NS}}}val", "") if style is not None else ""
    heading = re.match(r"(?:Heading|标题)\s*([1-6])", style_name, re.IGNORECASE)
    is_list = element.find("./w:pPr/w:numPr", NS) is not None
    if heading and text:
        text = f"{'#' * (int(heading.group(1)) + 1)} {text}"
    elif is_list and text:
        text = f"- {text}"
    return "\n\n".join(value for value in [text, *images] if value)


def word_table(element: ET.Element, relationships: dict[str, str], package: zipfile.ZipFile, stager: ImageStager, location: str) -> str:
    rows: list[list[object]] = []
    for row_index, row in enumerate(element.findall("./w:tr", NS), 1):
        cells: list[object] = []
        for column_index, cell in enumerate(row.findall("./w:tc", NS), 1):
            parts = [
                word_paragraph(paragraph, relationships, package, stager, f"{location}·表格 R{row_index}C{column_index}")
                for paragraph in cell.findall("./w:p", NS)
            ]
            cells.append("<br>".join(value.replace("\n", "<br>") for value in parts if value))
        rows.append(cells)
    return markdown_table(rows)


def convert_word(source: Path, stager: ImageStager) -> Conversion:
    sections: list[str] = []
    warnings: list[str] = []
    with zipfile.ZipFile(source) as package:
        document_name = "word/document.xml"
        root = ET.fromstring(package.read(document_name))
        relationships = relationship_map(package, document_name)
        body = root.find("./w:body", NS)
        if body is None:
            raise ConversionError("Word 文档缺少 document body。")
        paragraph_number = 0
        table_number = 0
        for child in body:
            if child.tag == f"{{{W_NS}}}p":
                paragraph_number += 1
                value = word_paragraph(child, relationships, package, stager, f"正文段落 {paragraph_number}")
            elif child.tag == f"{{{W_NS}}}tbl":
                table_number += 1
                value = word_table(child, relationships, package, stager, f"正文表格 {table_number}")
            else:
                value = ""
            if value:
                sections.append(value)
        for label, pattern in (("页眉", "word/header"), ("页脚", "word/footer"), ("脚注", "word/footnotes.xml"), ("尾注", "word/endnotes.xml"), ("批注", "word/comments.xml")):
            part_names = sorted(name for name in package.namelist() if name.startswith(pattern) and name.endswith(".xml"))
            part_text: list[str] = []
            for part_name in part_names:
                part_root = ET.fromstring(package.read(part_name))
                rels = relationship_map(package, part_name)
                for index, paragraph in enumerate(part_root.findall(".//w:p", NS), 1):
                    value = word_paragraph(paragraph, rels, package, stager, f"{label} {index}")
                    if value:
                        part_text.append(value)
            if part_text:
                sections.extend([f"## {label}", *part_text])
        document_xml = package.read(document_name)
        if b"txbxContent" in document_xml:
            warnings.append("文档含文本框；已提取可访问文字，但复杂文本框阅读顺序可能与版面不同。")
        if b"oMath" in document_xml:
            warnings.append("文档含 Office 公式对象；已尽量保留其中可访问文字，公式版式可能不完整。")
        if b"trackedChanges" in document_xml or b"<w:ins" in document_xml or b"<w:del" in document_xml:
            warnings.append("文档可能含修订记录；Markdown 不保留修订视图和接受/拒绝状态。")
    return Conversion("\n\n".join(sections), warnings, {"paragraphs": paragraph_number, "tables": table_number})


def iter_powerpoint_shapes(shapes: Iterable) -> list:
    flattened = []
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            flattened.extend(iter_powerpoint_shapes(shape.shapes))
        else:
            flattened.append(shape)
    return sorted(flattened, key=lambda item: (int(getattr(item, "top", 0)), int(getattr(item, "left", 0))))


def convert_powerpoint(source: Path, stager: ImageStager) -> Conversion:
    if not module_available("pptx"):
        raise DependencyError("PowerPoint 转换需要 python-pptx；请按 scripts/requirements.txt 安装。")
    from pptx import Presentation
    from pptx.shapes.picture import Picture

    presentation = Presentation(source)
    slides: list[str] = []
    warnings: list[str] = []
    images_unavailable = 0
    for slide_number, slide in enumerate(presentation.slides, 1):
        title_text = slide.shapes.title.text.strip() if slide.shapes.title is not None else ""
        content = [f"## 幻灯片 {slide_number}{f'：{title_text}' if title_text else ''}"]
        for shape_index, shape in enumerate(iter_powerpoint_shapes(slide.shapes), 1):
            if shape is slide.shapes.title:
                continue
            location = f"幻灯片 {slide_number}·对象 {shape_index}"
            # Picture also covers filled picture placeholders. `hasattr(image)`
            # evaluates the property and can itself raise "no embedded image".
            if isinstance(shape, Picture):
                try:
                    image = shape.image
                    image_data, image_extension = image.blob, "." + image.ext
                except (ValueError, KeyError, AttributeError, OSError):
                    images_unavailable += 1
                    embedded = shape._element.xpath("./p:blipFill/a:blip/@r:embed")
                    linked = shape._element.xpath("./p:blipFill/a:blip/@r:link")
                    if linked and not embedded:
                        reason = "外部链接图片未内嵌到 PPT，未读取外部链接目标"
                    elif embedded:
                        reason = "内嵌图片引用失效或图片数据无法读取"
                    else:
                        reason = "图片对象没有可读取的内嵌图片引用"
                    warning = f"{location}：{reason}；请在原 PPT 中核对，必要时将图片嵌入后重新转换。"
                    warnings.append(warning)
                    content.append(f"> [!warning] 图片未能提取\n> {warning}")
                else:
                    # Only image decoding errors are recoverable here. Staging
                    # or output failures must still abort instead of losing data.
                    content.append(stager.add(image_data, image_extension, f"幻灯片-{slide_number}-图片-{shape_index}", location))
            elif getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                content.append(markdown_table(rows))
            elif getattr(shape, "has_text_frame", False):
                paragraphs: list[str] = []
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        paragraphs.append(f"{'  ' * paragraph.level}- {text}" if paragraph.level else text)
                if paragraphs:
                    content.append("\n".join(paragraphs))
            if getattr(shape, "has_chart", False):
                warnings.append(f"幻灯片 {slide_number} 含图表对象；已保留可访问文字，图表视觉样式可能不完整。")
        if slide.has_notes_slide:
            notes = getattr(slide.notes_slide, "notes_text_frame", None)
            notes_text = notes.text.strip() if notes is not None else ""
            if notes_text:
                content.extend(["### 讲者备注", notes_text])
        slides.append("\n\n".join(content))
    return Conversion("\n\n".join(slides), warnings, {"slides": len(presentation.slides), "images_unavailable": images_unavailable})


def pdf_text_block(block: dict) -> str:
    lines = []
    for line in block.get("lines", []):
        text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
        if text:
            lines.append(text)
    return "\n".join(lines)


def pdf_vector_bounds(page):
    """Union visible drawing bounds in unrotated, crop-relative coordinates."""
    import fitz

    visible = page.rect * page.derotation_matrix
    region = None
    for drawing in page.get_drawings():
        if drawing.get("rect") is None:
            continue
        rect = fitz.Rect(drawing["rect"])
        # Pad before testing emptiness: a horizontal/vertical path has zero
        # height/width. Stroke joins can extend beyond its centerline bounds.
        padding = 2.0 + abs(float(drawing.get("width") or 0)) * 5.0
        rect = fitz.Rect(rect.x0 - padding, rect.y0 - padding, rect.x1 + padding, rect.y1 + padding) & visible
        if not rect.is_empty:
            region = rect if region is None else region | rect
    return region


def pdf_vector_block(page, blocks: list[dict]) -> tuple[dict | None, set[int]]:
    """Include whole intersecting text/image blocks to avoid clipping labels."""
    import fitz

    region = pdf_vector_bounds(page)
    covered: set[int] = set()
    if region is None:
        return None, covered
    visible = page.rect * page.derotation_matrix
    # Expanding for one label can touch another block. Reach a fixed point so
    # every block covered by the rendered region is emitted exactly once.
    while True:
        previous_count = len(covered)
        for index, block in enumerate(blocks):
            if index in covered or "bbox" not in block:
                continue
            if not (block.get("type") == 0 and pdf_text_block(block)) and not (block.get("type") == 1 and block.get("image")):
                continue
            rect = fitz.Rect(block["bbox"]) & visible
            if rect.is_empty or not region.intersects(rect):
                continue
            covered.add(index)
            padded = fitz.Rect(rect.x0 - 1, rect.y0 - 1, rect.x1 + 1, rect.y1 + 1) & visible
            region = region | padded
        if len(covered) == previous_count:
            break
    return {"type": "vector", "bbox": tuple(region)}, covered


def convert_pdf(source: Path, stager: ImageStager) -> Conversion:
    if not module_available("fitz"):
        raise DependencyError("PDF 转换需要 PyMuPDF；请按 scripts/requirements.txt 安装。")
    import fitz

    pages: list[str] = []
    vector_pages: list[int] = []
    warnings = ["PDF 的多栏、复杂表格、公式和矢量图可能无法完全恢复原始版式。"]
    with fitz.open(source) as document:
        for page_number, page in enumerate(document, 1):
            content = [f"## 第 {page_number} 页"]
            blocks = page.get_text("dict").get("blocks", [])
            vector, covered = pdf_vector_block(page, blocks)
            if vector is not None:
                if stager.mode not in {"attachments", "multimodal"}:
                    raise ConversionError(f"PDF 第 {page_number} 页含矢量图，请选择 attachments 或 multimodal 模式。")
                blocks = [block for index, block in enumerate(blocks) if index not in covered] + [vector]
            blocks = sorted(blocks, key=lambda item: (item.get("bbox", [0, 0])[1], item.get("bbox", [0, 0])[0]))
            found_content = False
            for block_index, block in enumerate(blocks, 1):
                if block.get("type") == "vector":
                    # Text/drawing coordinates are unrotated; get_pixmap's clip
                    # uses rotated page coordinates. Preserve the displayed view.
                    clip = fitz.Rect(block["bbox"]) * page.rotation_matrix
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), clip=clip, alpha=False)
                    content.append(stager.add(pixmap.tobytes("png"), ".png", f"第-{page_number}-页-矢量图", f"第 {page_number} 页·矢量图区域（含区内文字及图片）"))
                    vector_pages.append(page_number)
                    warnings.append(f"第 {page_number} 页的矢量图已按所在区域渲染为图片；区内文字及位图随区域保留，不再单独重复输出。")
                    found_content = True
                elif block.get("type") == 0:
                    text = pdf_text_block(block)
                    if text:
                        content.append(text)
                        found_content = True
                elif block.get("type") == 1 and block.get("image"):
                    extension = "." + str(block.get("ext") or "png")
                    content.append(stager.add(block["image"], extension, f"第-{page_number}-页-图片-{block_index}", f"第 {page_number} 页·图片块 {block_index}"))
                    found_content = True
            if not found_content and stager.mode in {"attachments", "multimodal"}:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                content.append(stager.add(pixmap.tobytes("png"), ".png", f"第-{page_number}-页", f"第 {page_number} 页·整页扫描"))
            pages.append("\n\n".join(content))
        page_count = len(document)
    return Conversion("\n\n".join(pages), warnings, {"pages": page_count, "vector_pages": vector_pages})


def convert_source(source: Path, stager: ImageStager) -> Conversion:
    extension = source.suffix.lower()
    if extension == ".csv":
        return convert_delimited(source, ",")
    if extension == ".tsv":
        return convert_delimited(source, "\t")
    if extension in {".xlsx", ".xlsm"}:
        return convert_excel(source, stager)
    if extension == ".docx":
        return convert_word(source, stager)
    if extension == ".pptx":
        return convert_powerpoint(source, stager)
    if extension == ".pdf":
        return convert_pdf(source, stager)
    raise ConversionError(f"尚不支持格式：{extension}")


def build_note(source: Path, title: str, mode: str, source_hash: str, conversion: Conversion) -> str:
    now = datetime.now().astimezone()
    modified = datetime.fromtimestamp(source.stat().st_mtime, timezone.utc).astimezone()
    description = f"由 {source.name} 转换生成的完整来源文档；保留原始结构并记录不可转换项。"
    frontmatter = [
        "---",
        f"title: {yaml_string(title)}",
        f"created: {yaml_string(now.strftime('%Y-%m-%d %H:%M'))}",
        f"updated: {yaml_string(now.strftime('%Y-%m-%d %H:%M'))}",
        "source: file-conversion",
        f"source_file: {yaml_string(source.name)}",
        f"source_path: {yaml_string(str(source))}",
        f"source_format: {source.suffix.lower().lstrip('.')}",
        f"source_bytes: {source.stat().st_size}",
        f"source_modified: {yaml_string(modified.isoformat())}",
        f"source_sha256: {source_hash}",
        f"conversion_mode: {mode}",
        "type: source",
        "status: inbox",
        "description: >",
        f"  {description}",
        "aliases: []",
        "domain: []",
        "project:",
        "maturity: seed",
        "retrieval_priority: low",
        "route_to:",
        "tags:",
        "  - inbox",
        "  - file-conversion",
        "related: []",
        "---",
        "",
        f"# {title}",
        "",
        f"> [!info] 转换来源\n> `{source.name}` · {source.stat().st_size} bytes · `{mode}`",
        "",
        conversion.markdown.strip() or "[未提取到可用正文]",
    ]
    if conversion.warnings:
        frontmatter.extend(["", "## 转换说明", "", *[f"- {item}" for item in dict.fromkeys(conversion.warnings)]])
    return "\n".join(frontmatter).rstrip() + "\n"


def validate_source(path_value: str | None) -> Path:
    if not path_value:
        raise ConversionError("请提供源文件路径。")
    source = Path(path_value).expanduser().resolve()
    try:
        is_file = source.is_file()
    except OSError as error:
        raise source_access_error(source, error) from error
    if not is_file:
        raise SourceError(source, "source_not_found", f"源文件不存在或不是文件，请核对完整路径：{source}")
    extension = source.suffix.lower()
    if extension in LEGACY:
        raise ConversionError(f"旧格式 {extension} 不能可靠直接转换；请先用原应用另存为新版 OOXML 格式。")
    if extension not in SUPPORTED:
        raise ConversionError(f"不支持的格式 {extension}；支持：{', '.join(sorted(SUPPORTED))}")
    return source


def validate_vault(path_value: str | None) -> Path:
    if not path_value:
        raise ConversionError("转换时必须提供 --vault-root。")
    vault = Path(path_value).expanduser().resolve()
    required = [vault / "AGENTS.md", vault / "01_Inbox", vault / "07_Attachments"]
    if not all(path.exists() for path in required):
        raise ConversionError(f"目标不是已初始化的 Knowledge Vault：{vault}")
    return vault


def run_conversion(source: Path, vault: Path, mode: str, title_value: str | None) -> dict:
    inspection = inspect_source(source)
    if inspection["missing_dependencies"]:
        raise DependencyError("缺少转换依赖：" + ", ".join(inspection["missing_dependencies"]))
    if inspection["requires_image_mode"] and mode == "text":
        raise ConversionError("检测到图文混排；请先让用户选择 --mode attachments 或 --mode multimodal。")
    if source.suffix.lower() in {".xlsx", ".xlsm"} and int(inspection["image_count"] or 0) > 0 and mode == "text":
        mode = "attachments"
    title = safe_name(title_value or source.stem)
    attachment_name, attachment_dir = find_attachment_folder(vault, title)
    inbox_name = f"{datetime.now().strftime('%Y-%m-%d %H%M')} - {title}.md"
    inbox_path = unique_path(vault / "01_Inbox" / inbox_name)
    before_stat = source.stat()
    before_hash = inspection["sha256"]
    with tempfile.TemporaryDirectory(prefix="knowledge-capture-") as temp_value:
        stager = ImageStager(Path(temp_value), attachment_dir, attachment_name, mode)
        conversion = convert_source(source, stager)
        conversion.warnings.extend(stager.warnings)
        after_stat = source.stat()
        after_hash = sha256_file(source)
        if before_stat.st_size != after_stat.st_size or before_stat.st_mtime_ns != after_stat.st_mtime_ns or before_hash != after_hash:
            raise ConversionError("源文件在转换过程中发生变化，已停止写入。")
        return commit_conversion(source, title, mode, before_hash, inbox_path, conversion, stager)


def commit_conversion(
    source: Path,
    title: str,
    mode: str,
    source_hash: str,
    inbox_path: Path,
    conversion: Conversion,
    stager: ImageStager,
) -> dict:
    """Publish an already converted and source-verified document once."""
    note = build_note(source, title, mode, source_hash, conversion)
    copied: list[Path] = []
    temporary_note = inbox_path.with_suffix(inbox_path.suffix + ".tmp")
    try:
        if mode == "attachments" and stager.records:
            stager.destination.mkdir(parents=True, exist_ok=True)
            for record in stager.records:
                target = stager.destination / record.filename
                shutil.copy2(record.source, target)
                copied.append(target)
        temporary_note.write_text(note, encoding="utf-8", newline="\n")
        os.replace(temporary_note, inbox_path)
    except Exception:
        for path in copied:
            path.unlink(missing_ok=True)
        temporary_note.unlink(missing_ok=True)
        raise
    return {
        "status": "ok",
        "source": str(source),
        "source_sha256": source_hash,
        "mode": mode,
        "markdown": str(inbox_path),
        "attachment_directory": str(stager.destination) if copied else None,
        "images_processed": len(stager.records),
        "attachments_written": len(copied),
        "warnings": conversion.warnings,
        "details": conversion.details,
    }


def write_atomic(path: Path, content: str) -> None:
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(content, encoding="utf-8", newline="\n")
    os.replace(temporary, path)


def source_is_unchanged(source: Path, expected: dict) -> bool:
    stat = source.stat()
    return bool(
        stat.st_size == expected.get("bytes")
        and stat.st_mtime_ns == expected.get("modified_ns")
        and sha256_file(source) == expected.get("sha256")
    )


def prepare_multimodal(source: Path, vault: Path, title_value: str | None) -> dict:
    inspection = inspect_source(source)
    if inspection["missing_dependencies"]:
        raise DependencyError("缺少转换依赖：" + ", ".join(inspection["missing_dependencies"]))
    if source.suffix.lower() not in MIXED_LAYOUT or not inspection["requires_image_mode"]:
        fallback_mode = "attachments" if source.suffix.lower() in {".xlsx", ".xlsm"} and inspection["image_count"] else "text"
        return run_conversion(source, vault, fallback_mode, title_value)

    title = safe_name(title_value or source.stem)
    inbox_name = f"{datetime.now().strftime('%Y-%m-%d %H%M')} - {title}.md"
    inbox_path = unique_path(vault / "01_Inbox" / inbox_name)
    before = source.stat()
    source_state = {
        "path": str(source),
        "bytes": before.st_size,
        "modified_ns": before.st_mtime_ns,
        "sha256": inspection["sha256"],
    }
    job_root = Path(tempfile.mkdtemp(prefix="knowledge-capture-multimodal-"))
    try:
        image_root = job_root / "images"
        image_root.mkdir()
        stager = ImageStager(image_root, image_root, "", "multimodal")
        conversion = convert_source(source, stager)
        conversion.warnings.extend(stager.warnings)
        if not source_is_unchanged(source, source_state):
            raise ConversionError("源文件在多模态准备过程中发生变化，已停止处理。")
        if not stager.records:
            # Media parts can remain in the ZIP even when no picture can be
            # extracted. Re-entering text conversion would reject that same ZIP
            # as mixed layout and discard the useful per-object warnings.
            result = commit_conversion(source, title, "text", source_state["sha256"], inbox_path, conversion, stager)
            shutil.rmtree(job_root)
            return result

        draft_path = job_root / "draft.md"
        results_path = job_root / "results.json"
        manifest_path = job_root / "manifest.json"
        draft = build_note(source, title, "multimodal", source_state["sha256"], conversion)
        write_atomic(draft_path, draft)
        images = [
            {
                "id": record.identifier,
                "path": str(record.source),
                "sha256": sha256_file(record.source),
                "location": record.location,
                "model_readable": record.model_readable,
                "unreadable_reason": record.unreadable_reason,
            }
            for record in stager.records
        ]
        manifest = {
            "schema_version": MULTIMODAL_SCHEMA_VERSION,
            "source": source_state,
            "vault_root": str(vault),
            "title": title,
            "draft": str(draft_path),
            "planned_markdown": str(inbox_path),
            "results_file": str(results_path),
            "images": images,
        }
        write_atomic(manifest_path, json.dumps(manifest, ensure_ascii=False, indent=2) + "\n")
        return {
            "status": "needs_multimodal",
            "mode": "multimodal",
            "source": str(source),
            "source_sha256": source_state["sha256"],
            "manifest": str(manifest_path),
            "results_file": str(results_path),
            "images": images,
            "images_to_read": sum(1 for image in images if image["model_readable"]),
            "warnings": conversion.warnings,
            "details": conversion.details,
            "results_contract": {
                "images": [
                    {
                        "id": "I001",
                        "markdown": "忠实转写图片文字、表格、图表或示意关系；无法确认时明确标记",
                        "confidence": "high|medium|low",
                        "uncertainties": ["可选：无法确认的局部内容"],
                    }
                ]
            },
        }
    except Exception:
        if job_root.exists():
            shutil.rmtree(job_root, ignore_errors=True)
        raise


def load_json_object(path: Path, label: str) -> dict:
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as error:
        raise ConversionError(f"无法读取{label}：{error}") from error
    if not isinstance(value, dict):
        raise ConversionError(f"{label}顶层必须是 JSON 对象。")
    return value


def validate_multimodal_job_path(manifest_path: Path) -> Path:
    job_root = manifest_path.parent.resolve()
    temp_root = Path(tempfile.gettempdir()).resolve()
    if job_root.parent != temp_root or not job_root.name.startswith("knowledge-capture-multimodal-"):
        raise ConversionError("多模态 manifest 必须位于知识收创建的系统临时目录。")
    return job_root


def resolved_job_file(value: object, job_root: Path, label: str) -> Path:
    if not isinstance(value, str) or not value.strip():
        raise ConversionError(f"多模态 manifest 缺少 {label}。")
    path = Path(value).expanduser().resolve()
    if path.parent != job_root:
        raise ConversionError(f"多模态 manifest 中的 {label} 越出临时任务目录。")
    return path


def render_multimodal_result(image: dict, result: dict | None) -> str:
    location = str(image.get("location") or "位置未知")
    if not image.get("model_readable"):
        reason = str(image.get("unreadable_reason") or "图片格式不受模型读取工具支持")
        return f"> [!warning] 图片未完成多模态识别\n> 位置：{location}\n> 原因：{reason}"

    assert result is not None
    confidence = str(result["confidence"])
    uncertainties = result.get("uncertainties", [])
    header = [
        "> [!info] 图片多模态识别",
        f"> 位置：{location}",
        f"> 置信度：{confidence}",
    ]
    if uncertainties:
        header.append("> 不确定项：" + "；".join(uncertainties))
    if confidence == "low":
        header.append("> [!warning] 低置信度结果必须对照原文件复核。")
    return "\n".join(header) + "\n\n" + str(result["markdown"]).strip()


def apply_multimodal(manifest_value: str, results_value: str | None, vault: Path, cleanup: bool) -> dict:
    manifest_path = Path(manifest_value).expanduser().resolve()
    job_root = validate_multimodal_job_path(manifest_path)
    manifest = load_json_object(manifest_path, "多模态 manifest")
    if manifest.get("schema_version") != MULTIMODAL_SCHEMA_VERSION:
        raise ConversionError("多模态 manifest 版本不受支持。")
    if Path(str(manifest.get("vault_root") or "")).resolve() != vault:
        raise ConversionError("多模态 manifest 的 Vault 与当前 --vault-root 不一致。")

    source_data = manifest.get("source")
    if not isinstance(source_data, dict):
        raise ConversionError("多模态 manifest 缺少来源状态。")
    source = validate_source(str(source_data.get("path") or ""))
    if not source_is_unchanged(source, source_data):
        raise ConversionError("源文件在多模态识别期间发生变化，已停止写入。")

    draft_path = resolved_job_file(manifest.get("draft"), job_root, "draft")
    expected_results_path = resolved_job_file(manifest.get("results_file"), job_root, "results_file")
    results_path = Path(results_value).expanduser().resolve() if results_value else expected_results_path
    if results_path != expected_results_path:
        raise ConversionError("--results 必须使用 prepare 返回的 results_file。")
    results = load_json_object(results_path, "多模态识别结果")

    images = manifest.get("images")
    if not isinstance(images, list) or not images:
        raise ConversionError("多模态 manifest 没有图片清单。")
    expected_ids = {
        str(image.get("id"))
        for image in images
        if isinstance(image, dict) and image.get("model_readable")
    }
    result_items = results.get("images")
    if not isinstance(result_items, list):
        raise ConversionError("多模态识别结果必须包含 images 数组。")
    result_map: dict[str, dict] = {}
    errors: list[str] = []
    for index, item in enumerate(result_items, 1):
        if not isinstance(item, dict):
            errors.append(f"images[{index}] 必须是对象")
            continue
        identifier = str(item.get("id") or "")
        if not identifier or identifier in result_map:
            errors.append(f"images[{index}] 的 id 缺失或重复")
            continue
        markdown = item.get("markdown")
        confidence = item.get("confidence")
        uncertainties = item.get("uncertainties", [])
        if not isinstance(markdown, str) or not markdown.strip():
            errors.append(f"{identifier} 缺少非空 markdown")
        if confidence not in {"high", "medium", "low"}:
            errors.append(f"{identifier} 的 confidence 必须是 high、medium 或 low")
        if not isinstance(uncertainties, list) or not all(isinstance(value, str) and value.strip() for value in uncertainties):
            errors.append(f"{identifier} 的 uncertainties 必须是非空字符串数组")
        if isinstance(markdown, str) and MULTIMODAL_PLACEHOLDER in markdown:
            errors.append(f"{identifier} 的 markdown 含保留占位符")
        result_map[identifier] = item
    actual_ids = set(result_map)
    if actual_ids != expected_ids:
        errors.append(
            "图片结果 ID 不完整："
            f"missing={sorted(expected_ids - actual_ids)}, extra={sorted(actual_ids - expected_ids)}"
        )
    if errors:
        raise ConversionError("多模态识别结果校验失败：" + "；".join(errors))

    draft = draft_path.read_text(encoding="utf-8")
    for image in images:
        if not isinstance(image, dict) or not isinstance(image.get("id"), str):
            raise ConversionError("多模态 manifest 含无效图片记录。")
        image_path = resolved_job_file(image.get("path"), job_root / "images", f"{image['id']} path")
        if not image_path.is_file() or sha256_file(image_path) != image.get("sha256"):
            raise ConversionError(f"多模态临时图片缺失或已变化：{image['id']}")
        token = f"<!-- {MULTIMODAL_PLACEHOLDER}:{image['id']} -->"
        if draft.count(token) != 1:
            raise ConversionError(f"多模态草稿中的图片占位符异常：{image['id']}")
        draft = draft.replace(token, render_multimodal_result(image, result_map.get(image["id"])))
    if MULTIMODAL_PLACEHOLDER in draft:
        raise ConversionError("多模态草稿仍有未处理图片占位符。")
    if not source_is_unchanged(source, source_data):
        raise ConversionError("源文件在多模态结果应用前发生变化，已停止写入。")

    planned = Path(str(manifest.get("planned_markdown") or "")).expanduser().resolve()
    if planned.parent != (vault / "01_Inbox").resolve():
        raise ConversionError("多模态 manifest 的输出路径不在 01_Inbox。")
    output = unique_path(planned)
    write_atomic(output, draft)
    cleanup_warning = None
    if cleanup:
        try:
            shutil.rmtree(job_root)
        except OSError as error:
            cleanup_warning = str(error)
    return {
        "status": "ok",
        "source": str(source),
        "source_sha256": source_data["sha256"],
        "mode": "multimodal",
        "markdown": str(output),
        "images_processed": len(images),
        "images_recognized": len(expected_ids),
        "images_unreadable": len(images) - len(expected_ids),
        "temporary_files_removed": cleanup and cleanup_warning is None,
        "cleanup_warning": cleanup_warning,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source", nargs="?", help="Excel/CSV/PDF/DOCX/PPTX source file")
    parser.add_argument("--inspect", action="store_true", help="Inspect format, images, hash, and dependencies without writing")
    parser.add_argument("--vault-root", help="Initialized Knowledge Vault root used for final output")
    parser.add_argument("--mode", choices=("text", "attachments", "multimodal"), default="text")
    parser.add_argument("--title", help="Override the generated Inbox note title")
    parser.add_argument("--apply-multimodal", metavar="MANIFEST", help="Apply model-produced image results to a prepared draft")
    parser.add_argument("--results", help="Model-produced multimodal results JSON returned by prepare")
    parser.add_argument("--cleanup", action="store_true", help="Remove temporary multimodal files after verified success")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        if args.apply_multimodal and args.inspect:
            raise ConversionError("--inspect 不能与 --apply-multimodal 同时使用。")
        if args.results and not args.apply_multimodal:
            raise ConversionError("--results 只能与 --apply-multimodal 同时使用。")
        if args.cleanup and not args.apply_multimodal:
            raise ConversionError("--cleanup 只能与 --apply-multimodal 同时使用。")
        if args.apply_multimodal and args.source:
            raise ConversionError("应用多模态结果时不再传入源文件位置参数。")
        if args.apply_multimodal:
            result = apply_multimodal(args.apply_multimodal, args.results, validate_vault(args.vault_root), args.cleanup)
        else:
            source = validate_source(args.source)
            if args.inspect:
                result = inspect_source(source)
            elif args.mode == "multimodal":
                result = prepare_multimodal(source, validate_vault(args.vault_root), args.title)
            else:
                result = run_conversion(source, validate_vault(args.vault_root), args.mode, args.title)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except SourceError as error:
        print(json.dumps({"status": "error", "error_code": error.error_code, "source": str(error.source), "error": str(error)}, ensure_ascii=False, indent=2), file=sys.stderr)
        return 2
    except (ConversionError, DependencyError, zipfile.BadZipFile) as error:
        print(json.dumps({"status": "error", "error": str(error)}, ensure_ascii=False, indent=2), file=sys.stderr)
        return 2
    except Exception as error:
        print(json.dumps({"status": "error", "error": f"转换失败：{error}"}, ensure_ascii=False, indent=2), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
